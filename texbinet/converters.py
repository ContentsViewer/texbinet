from pathlib import Path


def pdf2text(path: Path) -> str:
    """Convert a PDF file to text."""

    # import easyocr
    # from pdf2image import convert_from_path

    # images = convert_from_path(path)
    # reader = easyocr.Reader(["en", "ja"])

    # text = ""
    # for image in images:
    #     texts = reader.readtext(np.array(image), detail=0)
    #     print(texts)
    #     text += "\n".join(texts) + "\n"

    # return text

    import pypdf

    with path.open("rb") as f:
        pdf = pypdf.PdfReader(f)
        return "\n".join(page.extract_text() for page in pdf.pages)


def image2text(path: Path) -> str:
    """Convert an image file to text."""

    import easyocr

    reader = easyocr.Reader(["en", "ja"])
    texts = reader.readtext(str(path))

    text = "\n".join([text for _, text, accuracy in texts if accuracy > 0.5])

    return text

    # import torch
    # from transformers import Blip2Processor, Blip2ForConditionalGeneration
    # from PIL import Image

    # image = Image.open(path)

    # processor = Blip2Processor.from_pretrained("Salesforce/blip2-opt-2.7b")
    # model = Blip2ForConditionalGeneration.from_pretrained("Salesforce/blip2-opt-2.7b", load_in_8bit=True, device_map="auto")

    # device = "cuda" if torch.cuda.is_available() else "cpu"

    # # processor = Blip2Processor.from_pretrained("Salesforce/blip2-opt-2.7b")
    # # model = Blip2ForConditionalGeneration.from_pretrained(
    # #     "Salesforce/blip2-opt-2.7b", torch_dtype=torch.float16
    # # ).to(device)

    # inputs = processor(images=image, return_tensors="pt").to(device, torch.float16)
    # generated_ids = model.generate(**inputs, max_new_tokens=20)
    # # generated_text = processor.batch_decode(generated_ids, skip_special_tokens=True)[
    # #     0
    # # ].strip()
    # generated_text = processor.decode(generated_ids[0], skip_special_tokens=True).strip()
    # print(generated_text)
    # return generated_text


def pptx2text(path: Path) -> str:
    """Convert a PPTX file to text."""

    import pptx
    from typing import List

    def extract_table_text(table) -> List[str]:
        """Extract text from table shape in PPTX."""
        texts = []
        for row in table.rows:
            row_text = []
            for cell in row.cells:
                cell_text = "".join(
                    run.text
                    for paragraph in cell.text_frame.paragraphs
                    for run in paragraph.runs
                )
                row_text.append(cell_text.replace("\n", ""))
            texts.append(", ".join(row_text))
        return texts

    prs = pptx.Presentation(path)

    texts = []

    for i, slide in enumerate(prs.slides, start=1):
        texts.append(f"--- Page {i} ---")

        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text.replace("\n", ""))

            if shape.has_table:
                texts.extend(extract_table_text(shape.table))

    return "\n".join(texts)
