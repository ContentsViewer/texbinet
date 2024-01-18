from pathlib import Path
from typing import List

import pptx


def pptx2text(path: Path) -> str:
    """Convert a PPTX file to text."""

    def extract_table_text(table) -> List[str]:
        """Extract text from table shape in PPTX."""
        texts = []
        for row in table.rows:
            row_text = []
            for cell in row.cells:
                cell_text = "".join(
                    run.text
                    for paragraph in cell.text_frame.paragraphs
                    for run in paragraph.runs
                )
                row_text.append(cell_text.replace("\n", ""))
            texts.append(", ".join(row_text))
        return texts

    prs = pptx.Presentation(path)

    texts = []

    for i, slide in enumerate(prs.slides, start=1):
        texts.append(f"--- Page {i} ---")

        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text.replace("\n", ""))

            if shape.has_table:
                texts.extend(extract_table_text(shape.table))

    return "\n".join(texts)
